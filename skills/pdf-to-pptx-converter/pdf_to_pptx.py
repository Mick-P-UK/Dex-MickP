#!/usr/bin/env python3
"""
PDF to PPTX Converter - DIY Investors Branded
Version: 2.1
Change: Fixed colour sampling region. Previously sampled within the NotebookLM logo
        bar itself, producing a grey/dark mismatch. Now samples 120x20px strip
        positioned 50-70px ABOVE the bottom edge, in clear slide background.
        Added built-in post-conversion verification reporting sampled colours per slide.

Usage:
    python3 pdf_to_pptx.py <input.pdf> <output.pptx> <logo.jpg> [event_label]
"""

import sys
import os
import statistics
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image
import io
import tempfile

# ---------------------------------------------------------------------------
# Slide dimensions (16:9)
# ---------------------------------------------------------------------------
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)

# ---------------------------------------------------------------------------
# Branding header
# ---------------------------------------------------------------------------
HEADER_H     = Inches(0.34)
LOGO_W       = Inches(1.2)
LOGO_H       = Inches(0.24)
BANNER_COLOR = RGBColor(192, 57, 43)   # DIY Investors red

# ---------------------------------------------------------------------------
# NotebookLM masking rectangle (bottom-right corner, absolute)
# ---------------------------------------------------------------------------
MASK_W = Inches(1.20)
MASK_H = Inches(0.17)
MASK_X = SLIDE_W - MASK_W   # right edge
MASK_Y = SLIDE_H - MASK_H   # bottom edge

# ---------------------------------------------------------------------------
# Colour sampling
# FIX v2.1: Sample ABOVE the NotebookLM logo strip, not within it.
# The logo bar occupies roughly the bottom 40px of the image. Sampling
# within that region picks up the logo's own dark colour, not the slide
# background. Sample a 120x20px patch 50-70px above the bottom edge.
# ---------------------------------------------------------------------------
SAMPLE_W_PX      = 120
SAMPLE_H_PX      = 20
SAMPLE_OFFSET_PX = 70   # px from bottom -- clears the logo bar with margin

PDF_DPI     = 200
BORDER_CROP = 20   # px to crop on each edge before processing


# ---------------------------------------------------------------------------
# Warning heuristic
# The NotebookLM logo bar is a specific dark grey: approx RGB(85-110, 90-115, 95-120).
# We flag only colours that fall squarely in that narrow band as likely wrong samples.
# ---------------------------------------------------------------------------
def _looks_like_logo_bar(r, g, b):
    """Return True if the sampled colour looks like the NotebookLM logo bar itself."""
    return all(80 <= v <= 125 for v in (r, g, b))


# ---------------------------------------------------------------------------
# Core helpers
# ---------------------------------------------------------------------------

def extract_page_image(page, dpi=PDF_DPI) -> Image.Image:
    """Render a PDF page to a PIL Image at the given DPI."""
    mat = fitz.Matrix(dpi / 72, dpi / 72)
    pix = page.get_pixmap(matrix=mat, alpha=False)
    img_bytes = pix.tobytes("png")
    return Image.open(io.BytesIO(img_bytes)).convert("RGB")


def crop_borders(img: Image.Image, border: int = BORDER_CROP) -> Image.Image:
    """Crop a uniform border from all four sides."""
    w, h = img.size
    return img.crop((border, border, w - border, h - border))


def sample_background_color(img: Image.Image) -> RGBColor:
    """
    Sample the slide background colour from just ABOVE the NotebookLM logo.

    The NotebookLM watermark sits in the bottom ~40px of each slide image.
    Sampling within the logo picks up its own dark/grey bar colour, NOT the
    slide background. Instead we sample a 120x20px strip positioned
    SAMPLE_OFFSET_PX px above the bottom edge -- well clear of the logo.

    Median is used rather than mean to be robust against dark gradients or
    UI elements that may appear near the edge.
    """
    iw, ih = img.size
    sx = max(0, iw - SAMPLE_W_PX - 10)   # near right edge, slight inset
    sy = max(0, ih - SAMPLE_OFFSET_PX)   # above the logo strip
    region = img.crop((sx, sy, sx + SAMPLE_W_PX, sy + SAMPLE_H_PX))
    pixels = list(region.convert("RGB").getdata())
    if not pixels:
        return RGBColor(255, 255, 255)
    r = int(statistics.median(p[0] for p in pixels))
    g = int(statistics.median(p[1] for p in pixels))
    b = int(statistics.median(p[2] for p in pixels))
    return RGBColor(r, g, b)


# ---------------------------------------------------------------------------
# Slide building
# ---------------------------------------------------------------------------

def add_header(slide, logo_path: str, event_label: str):
    """Add the DIY Investors branded header bar to a slide."""
    # Red banner across full width
    banner = slide.shapes.add_shape(1, 0, 0, SLIDE_W, HEADER_H)
    banner.fill.solid()
    banner.fill.fore_color.rgb = BANNER_COLOR
    banner.line.fill.background()

    # Logo (left)
    if os.path.exists(logo_path):
        logo_top = (HEADER_H - LOGO_H) // 2
        slide.shapes.add_picture(logo_path, Inches(0.05), logo_top, LOGO_W, LOGO_H)

    # Event label text (right-aligned)
    text_left  = LOGO_W + Inches(0.15)
    text_width = SLIDE_W - text_left - Inches(0.1)
    txBox = slide.shapes.add_textbox(text_left, 0, text_width, HEADER_H)
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = event_label
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)


def add_mask(slide, bg_color: RGBColor):
    """
    Add an invisible rectangle over the NotebookLM logo (bottom-right corner).
    Fill is colour-matched to the slide background so it is seamless.
    No border, no shadow.
    """
    mask = slide.shapes.add_shape(1, MASK_X, MASK_Y, MASK_W, MASK_H)
    mask.fill.solid()
    mask.fill.fore_color.rgb = bg_color
    mask.line.fill.background()
    # Disable shadow
    from pptx.oxml.ns import qn
    sp_pr = mask._element.spPr
    effect_lst = sp_pr.find(qn('a:effectLst'))
    if effect_lst is not None:
        sp_pr.remove(effect_lst)


# ---------------------------------------------------------------------------
# Main conversion
# ---------------------------------------------------------------------------

def convert_pdf_to_pptx(pdf_path: str, output_path: str, logo_path: str,
                         event_label: str = "DIY Investors"):
    """Convert a PDF slide deck to a branded PPTX. Returns per-slide colour log."""
    print(f"\nDIY Investors PDF -> PPTX Converter v2.1")
    print(f"Input : {pdf_path}")
    print(f"Output: {output_path}")
    print(f"Label : {event_label}")
    print("-" * 60)

    doc = fitz.open(pdf_path)
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]
    colour_log = []   # for verification report

    for page_num, page in enumerate(doc):
        slide_n = page_num + 1
        print(f"  Slide {slide_n:2d}/{len(doc)} ...", end=" ", flush=True)

        # 1. Render + crop
        img = extract_page_image(page)
        img = crop_borders(img)

        # 2. Sample background colour (ABOVE the logo, not within it)
        bg_color = sample_background_color(img)
        colour_log.append((slide_n, bg_color[0], bg_color[1], bg_color[2]))
        print(f"bg=RGB({bg_color[0]:3d},{bg_color[1]:3d},{bg_color[2]:3d})", end=" ")

        # 3. Save image to temp PNG
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
            tmp_path = tmp.name
            img.save(tmp_path, format="PNG", quality=95)

        # 4. Build slide: full-bleed image (below header)
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(tmp_path, 0, HEADER_H, SLIDE_W, SLIDE_H - HEADER_H)
        os.unlink(tmp_path)

        # 5. Branding header
        add_header(slide, logo_path, event_label)

        # 6. NotebookLM mask
        add_mask(slide, bg_color)

        print("OK")

    doc.close()
    prs.save(output_path)

    # -----------------------------------------------------------------------
    # Built-in verification report
    # -----------------------------------------------------------------------
    print("\n" + "=" * 60)
    print("VERIFICATION REPORT")
    print("=" * 60)
    print(f"{'Slide':<7} {'Sampled Background RGB':<28} {'Status'}")
    print("-" * 60)
    warn_count = 0
    for slide_n, r, g, b in colour_log:
        # Warn only if sampled colour falls in the NotebookLM logo bar's
        # characteristic dark grey band (RGB ~85-125 across all channels).
        # Note: legitimate dark backgrounds (e.g. navy) will have one channel
        # much lower than the others and will NOT trigger this warning.
        bad = _looks_like_logo_bar(r, g, b)
        if bad:
            warn_count += 1
        status = "WARN: may have sampled logo bar" if bad else "OK"
        print(f"  {slide_n:<5} RGB({r:3d}, {g:3d}, {b:3d})              {status}")

    print("-" * 60)
    print(f"Total slides : {len(colour_log)}")
    print(f"OK           : {len(colour_log) - warn_count}")
    print(f"Warnings     : {warn_count}")
    if warn_count == 0:
        print("Result       : PASS -- all masks colour-matched correctly")
    else:
        print("Result       : REVIEW slides marked WARN above")
    print("=" * 60)
    print(f"\nOutput saved: {output_path}\n")


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    if len(sys.argv) < 4:
        print("Usage: pdf_to_pptx.py <input.pdf> <output.pptx> <logo.jpg> [event_label]")
        sys.exit(1)
    label = sys.argv[4] if len(sys.argv) > 4 else "DIY Investors"
    convert_pdf_to_pptx(sys.argv[1], sys.argv[2], sys.argv[3], label)
