"""Helper library for building fully editable PowerPoint graphics.

Everything produced here is a native p:sp - freeform, preset autoshape or text
box - so Mick can move, recolour, resize and retype it in PowerPoint.

Work in INCHES with maths-convention angles (0 = east, counter-clockwise) and let
these helpers do the conversion to EMU and to OOXML's clockwise-from-3-o'clock
angle convention.
"""
import math

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

EMU = 914400
SLIDE_16x9 = (13.333, 7.5)


# ------------------------------------------------------------------ basics --
def E(inches):
    """Inches -> Emu."""
    return Emu(int(round(inches * EMU)))


def rgb(hex6):
    """'RRGGBB' (no hash) -> RGBColor."""
    return RGBColor.from_string(hex6.lstrip("#").upper())


def deck(width=SLIDE_16x9[0], height=SLIDE_16x9[1]):
    """New presentation with a single blank slide. Returns (prs, slide)."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = E(width), E(height)
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def polar(cx, cy, ang_deg, r):
    """Maths-convention polar -> slide coords in inches (y increases downward)."""
    a = math.radians(ang_deg)
    return (cx + r * math.cos(a), cy - r * math.sin(a))


# ------------------------------------------------------------------ styling --
def style(shape, fill, line=None, line_pt=0.0, name=None):
    """Solid fill, optional outline, no shadow, optional shape name.

    IMPORTANT: python-pptx writes a <p:style> block whose effectRef idx="2"
    re-introduces a drop shadow even after shadow.inherit = False. Renderers
    honour it. Zeroing lnRef/fillRef/effectRef is the actual fix.
    """
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(line_pt)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    st = shape._element.find(qn("p:style"))
    if st is not None:
        for tag in ("a:lnRef", "a:fillRef", "a:effectRef"):
            ref = st.find(qn(tag))
            if ref is not None:
                ref.set("idx", "0")
    if name:
        shape._element.nvSpPr.cNvPr.set("name", name)
    return shape


def set_adj(shape, **vals):
    """Write preset-geometry adjustment values directly, in raw OOXML units.

    Do not use shape.adjustments for angles - python-pptx divides by 100000,
    which is correct for ratios and wrong for the 60000ths-of-a-degree angles
    that blockArc and friends use.
    """
    av = shape._element.spPr.find(qn("a:prstGeom")).find(qn("a:avLst"))
    for gd in list(av):
        av.remove(gd)
    for k, v in vals.items():
        av.append(av.makeelement(qn("a:gd"), {"name": k, "fmla": "val %d" % v}))
    return shape


# ------------------------------------------------------------------- shapes --
def rect(slide, x, y, w, h, fill, line=None, line_pt=0.0, name=None, rounded=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        E(x), E(y), E(w), E(h))
    return style(shp, fill, line, line_pt, name)


def donut(slide, cx, cy, r_in, r_out, fill, name=None):
    """Annulus as a native DONUT autoshape (thickness stays adjustable in PPT)."""
    d = 2 * r_out
    shp = slide.shapes.add_shape(MSO_SHAPE.DONUT, E(cx - r_out), E(cy - r_out),
                                 E(d), E(d))
    set_adj(shp, adj=int(round((r_out - r_in) / d * 100000)))
    return style(shp, fill, name=name)


def block_arc(slide, cx, cy, r_in, r_out, centre_ang, span, fill, name=None):
    """Annular sector as a native BLOCK_ARC autoshape.

    centre_ang and span are in maths convention (degrees, counter-clockwise).
    OOXML measures clockwise from three o'clock in 60000ths of a degree and
    draws clockwise from adj1 to adj2, hence the negation.
    """
    side = 2 * r_out
    shp = slide.shapes.add_shape(MSO_SHAPE.BLOCK_ARC, E(cx - r_out), E(cy - r_out),
                                 E(side), E(side))
    set_adj(shp,
            adj1=int(round((-(centre_ang + span / 2.0)) % 360.0 * 60000)),
            adj2=int(round((-(centre_ang - span / 2.0)) % 360.0 * 60000)),
            adj3=int(round((r_out - r_in) / r_out * 100000)))
    return style(shp, fill, name=name)


def freeform(slide, vertices_in, fill, line=None, line_pt=0.0, name=None):
    """Closed freeform from a list of (x, y) inch pairs."""
    v = [(int(round(x * EMU)), int(round(y * EMU))) for x, y in vertices_in]
    fb = slide.shapes.build_freeform(v[0][0], v[0][1], scale=1.0)
    fb.add_line_segments(v[1:], close=True)
    return style(fb.convert_to_shape(), fill, line, line_pt, name)


def annulus_sector(cx, cy, r_in, r_out, a0, a1, n=None):
    """Vertices for an annular sector, for when a freeform beats a BLOCK_ARC
    (e.g. you need a non-circular or clipped ring)."""
    span = (a1 - a0) % 360.0
    n = n or max(24, int(span / 2.0))
    outer = [polar(cx, cy, a0 + span * i / n, r_out) for i in range(n + 1)]
    inner = [polar(cx, cy, a1 - span * i / n, r_in) for i in range(n + 1)]
    return outer + inner


# --------------------------------------------------------------------- text --
def textbox(slide, text, cx, cy, w, h, size, colour, rot=0.0, bold=True,
            italic=False, font="Calibri", name=None, spacing=None,
            align=PP_ALIGN.CENTER, wrap=False):
    """Centred, zero-margin text box positioned by its CENTRE point."""
    box = slide.shapes.add_textbox(E(cx - w / 2.0), E(cy - h / 2.0), E(w), E(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size, run.font.bold, run.font.italic = Pt(size), bold, italic
    run.font.name = font
    run.font.color.rgb = rgb(colour)
    if spacing:                       # character spacing in points
        run.font._rPr.set("spc", str(int(spacing * 100)))
    if rot:
        box.rotation = rot % 360
    if name:
        box._element.nvSpPr.cNvPr.set("name", name)
    return box


def ring_rotation(theta):
    """Rotation for a text box so a tangential ring label reads upright.

    Upper half: 'up' points outward. Lower half: 'up' points inward. Both look
    correct to the viewer. Keeps the label a normal, editable text box instead
    of un-editable WordArt on a path.
    """
    return (90.0 - theta) if math.sin(math.radians(theta)) >= 0 else (270.0 - theta)


def ring_text_room(r_label, r_out):
    """Horizontal inches available for a NON-rotated label sitting at radius
    r_label inside a ring of outer radius r_out. If your text is wider than
    this, shrink it or rotate it with ring_rotation()."""
    if r_label >= r_out:
        return 0.0
    return 2.0 * math.sqrt(r_out ** 2 - r_label ** 2)


# ------------------------------------------------------------------- jigsaw --
def jigsaw_boundary(cx, cy, ang, r_disc, knob_d=None, knob_r=None, knob_off=None,
                    n=28):
    """Points from disc centre out to the rim along `ang`, with a jigsaw knob.

    The knob always bulges toward INCREASING angle. Build each piece by walking
    its start boundary outward and its end boundary reversed, and every piece
    ends up with exactly one tab and one matching socket - neighbours always
    mesh, with no separate tab/socket bookkeeping.
    """
    knob_d = knob_d if knob_d is not None else 0.60 * r_disc
    knob_r = knob_r if knob_r is not None else 0.135 * r_disc
    knob_off = knob_off if knob_off is not None else 0.55 * knob_r

    a = math.radians(ang)
    u, p = (math.cos(a), math.sin(a)), (-math.sin(a), math.cos(a))
    w = math.sqrt(knob_r ** 2 - knob_off ** 2)          # half-width of the neck
    kc = (knob_d * u[0] + knob_off * p[0], knob_d * u[1] + knob_off * p[1])

    phi_a = math.degrees(math.atan2(-knob_off, -w))
    phi_b = math.degrees(math.atan2(-knob_off, w))
    if phi_a < phi_b:                                   # go the long way, via +p
        phi_a += 360.0

    pts = [(cx, cy), (cx + (knob_d - w) * u[0], cy - (knob_d - w) * u[1])]
    for i in range(n + 1):
        phi = math.radians(phi_a + (phi_b - phi_a) * i / n)
        pts.append((cx + kc[0] + knob_r * (math.cos(phi) * u[0] + math.sin(phi) * p[0]),
                    cy - kc[1] - knob_r * (math.cos(phi) * u[1] + math.sin(phi) * p[1])))
    pts.append(polar(cx, cy, ang, r_disc))
    return pts


def jigsaw_piece(cx, cy, a0, a1, r_disc, **knob):
    """Vertices for one interlocking piece spanning a0 -> a1 counter-clockwise."""
    span = (a1 - a0) % 360.0
    n = max(24, int(span / 2.0))
    rim = [polar(cx, cy, a0 + span * i / n, r_disc) for i in range(1, n + 1)]
    return (jigsaw_boundary(cx, cy, a0, r_disc, **knob) + rim
            + list(reversed(jigsaw_boundary(cx, cy, a1, r_disc, **knob))))
